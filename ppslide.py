from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

# Create presentation
prs = Presentation()
slide_layout = prs.slide_layouts[6]  # Blank slide
slide = prs.slides.add_slide(slide_layout)

# Define colors and styles
colors = {
    "Base 0": RGBColor(0, 102, 204),  # Blue
    "Base 1": RGBColor(0, 153, 255),
    "Base 2": RGBColor(51, 204, 255),
    "CoT": RGBColor(102, 204, 255)
}
title_font_size = Pt(32)
box_font_size = Pt(20)
count_font_size = Pt(28)

# Add title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
title_tf = title_box.text_frame
title_run = title_tf.paragraphs[0].add_run()
title_run.text = "Prompting Strategies"
title_run.font.size = title_font_size
title_run.font.bold = True

# Add GPT-3 source
gpt_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(0.5))
gpt_tf = gpt_box.text_frame
gpt_tf.text = "Model: GPT-3 (OpenAI)"
gpt_tf.paragraphs[0].runs[0].font.size = Pt(18)

# Add 4 prompt variant boxes horizontally
labels = ["Base 0: Zero-shot", "Base 1: One-shot", "Base 2: Few-shot", "CoT: Step-by-step"]
left_start = 0.5
top = 2
width = 2
height = 1.2

for i, label in enumerate(labels):
    left = Inches(left_start + i * (width + 0.2))
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = list(colors.values())[i]
    shape.line.color.rgb = RGBColor(255, 255, 255)

    text_frame = shape.text_frame
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = label
    run.font.size = box_font_size
    run.font.bold = True
    run.font.color.rgb = RGBColor(255, 255, 255)

# Add total justifications count
summary_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(1))
summary_tf = summary_box.text_frame
summary_run = summary_tf.paragraphs[0].add_run()
summary_run.text = "Each scenario prompted 4 times → 396 justifications"
summary_run.font.size = count_font_size
summary_run.font.bold = True

# Save presentation
prs.save("Prompting_Strategies_Slide.pptx")
